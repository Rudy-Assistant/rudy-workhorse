"""
OCR & Document Intelligence — Extract text from images, PDFs, documents.

Capabilities:
  - Image OCR: Extract text from screenshots, photos, scanned docs
  - PDF extraction: Text, tables, and structure from PDFs
  - Document parsing: DOCX, RTF, EPUB, PPTX text extraction
  - Batch processing: Process entire directories of documents
  - Multi-language: Support for 60+ languages via EasyOCR
  - Table detection: Extract tabular data from images/PDFs
"""

import json
import os

from pathlib import Path
from typing import List

DESKTOP = Path(os.environ.get("USERPROFILE", os.path.expanduser("~"))) / "Desktop"
OCR_DIR = DESKTOP / "rudy-data" / "ocr-output"
OCR_DIR.mkdir(parents=True, exist_ok=True)

class ImageOCR:
    """OCR for images using EasyOCR."""

    def __init__(self, languages: List[str] = None):
        self.languages = languages or ["en"]
        self._reader = None

    def _get_reader(self):
        if self._reader is None:
            try:
                import easyocr
                self._reader = easyocr.Reader(self.languages, gpu=False)
            except ImportError:
                raise RuntimeError("easyocr not installed. Run: pip install easyocr")
        return self._reader

    def read_image(self, image_path: str, detail: bool = True) -> dict:
        """
        Extract text from an image.
        Returns: {"text": str, "blocks": [...], "confidence": float}
        """
        reader = self._get_reader()
        results = reader.readtext(image_path)

        blocks = []
        full_text = []
        total_conf = 0

        for (bbox, text, conf) in results:
            blocks.append({
                "text": text,
                "confidence": round(conf, 3),
                "bbox": [[int(p[0]), int(p[1])] for p in bbox],
            })
            full_text.append(text)
            total_conf += conf

        avg_conf = total_conf / len(results) if results else 0

        return {
            "text": " ".join(full_text),
            "blocks": blocks if detail else [],
            "confidence": round(avg_conf, 3),
            "block_count": len(results),
            "source": image_path,
        }

    def read_batch(self, image_paths: List[str]) -> List[dict]:
        """OCR multiple images."""
        return [self.read_image(p) for p in image_paths]

class PDFExtractor:
    """Extract text and tables from PDFs."""

    def extract_text(self, pdf_path: str) -> dict:
        """Extract all text from a PDF."""
        try:
            import pdfplumber
            pages = []
            full_text = []

            with pdfplumber.open(pdf_path) as pdf:
                for i, page in enumerate(pdf.pages):
                    text = page.extract_text() or ""
                    pages.append({"page": i + 1, "text": text})
                    full_text.append(text)

            return {
                "text": "\n\n".join(full_text),
                "pages": pages,
                "page_count": len(pages),
                "source": pdf_path,
            }
        except ImportError:
            # Fallback to PyPDF2
            try:
                import PyPDF2
                pages = []
                full_text = []
                with open(pdf_path, "rb") as f:
                    reader = PyPDF2.PdfReader(f)
                    for i, page in enumerate(reader.pages):
                        text = page.extract_text() or ""
                        pages.append({"page": i + 1, "text": text})
                        full_text.append(text)
                return {
                    "text": "\n\n".join(full_text),
                    "pages": pages,
                    "page_count": len(pages),
                    "source": pdf_path,
                }
            except ImportError:
                raise RuntimeError("pdfplumber or PyPDF2 needed")

    def extract_tables(self, pdf_path: str) -> List[dict]:
        """Extract tables from a PDF."""
        try:
            import pdfplumber
            tables = []
            with pdfplumber.open(pdf_path) as pdf:
                for i, page in enumerate(pdf.pages):
                    page_tables = page.extract_tables()
                    for j, table in enumerate(page_tables):
                        if table:
                            tables.append({
                                "page": i + 1,
                                "table_index": j,
                                "headers": table[0] if table else [],
                                "rows": table[1:] if len(table) > 1 else [],
                                "row_count": len(table) - 1,
                            })
            return tables
        except ImportError:
            return [{"error": "pdfplumber not installed"}]

class DocumentParser:
    """Parse various document formats."""

    def parse(self, filepath: str) -> dict:
        """Auto-detect format and extract text."""
        path = Path(filepath)
        ext = path.suffix.lower()

        parsers = {
            ".pdf": self._parse_pdf,
            ".docx": self._parse_docx,
            ".doc": self._parse_docx,
            ".rtf": self._parse_rtf,
            ".epub": self._parse_epub,
            ".pptx": self._parse_pptx,
            ".txt": self._parse_text,
            ".md": self._parse_text,
            ".csv": self._parse_text,
            ".json": self._parse_json,
            ".html": self._parse_html,
            ".htm": self._parse_html,
        }

        parser = parsers.get(ext)
        if parser:
            return parser(filepath)
        return {"error": f"Unsupported format: {ext}", "source": filepath}

    def _parse_pdf(self, path):
        return PDFExtractor().extract_text(path)

    def _parse_docx(self, path):
        try:
            import docx2txt
            text = docx2txt.process(path)
            return {"text": text, "source": path, "format": "docx"}
        except ImportError:
            try:
                from docx import Document
                doc = Document(path)
                text = "\n".join(p.text for p in doc.paragraphs)
                return {"text": text, "source": path, "format": "docx"}
            except ImportError:
                return {"error": "docx2txt or python-docx needed"}

    def _parse_rtf(self, path):
        try:
            from striprtf.striprtf import rtf_to_text
            with open(path, encoding="utf-8", errors="replace") as f:
                text = rtf_to_text(f.read())
            return {"text": text, "source": path, "format": "rtf"}
        except ImportError:
            return {"error": "striprtf not installed"}

    def _parse_epub(self, path):
        try:
            import ebooklib
            from ebooklib import epub
            from bs4 import BeautifulSoup
            book = epub.read_epub(path)
            texts = []
            for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
                soup = BeautifulSoup(item.get_content(), "html.parser")
                texts.append(soup.get_text())
            return {"text": "\n\n".join(texts), "source": path, "format": "epub"}
        except ImportError:
            return {"error": "ebooklib not installed"}

    def _parse_pptx(self, path):
        try:
            from pptx import Presentation
            prs = Presentation(path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        texts.append(shape.text)
            return {"text": "\n\n".join(texts), "source": path, "format": "pptx"}
        except ImportError:
            return {"error": "python-pptx not installed"}

    def _parse_text(self, path):
        with open(path, encoding="utf-8", errors="replace") as f:
            return {"text": f.read(), "source": path, "format": "text"}

    def _parse_json(self, path):
        with open(path) as f:
            data = json.load(f)
        return {"text": json.dumps(data, indent=2), "source": path, "format": "json"}

    def _parse_html(self, path):
        try:
            from bs4 import BeautifulSoup
            with open(path, encoding="utf-8", errors="replace") as f:
                soup = BeautifulSoup(f.read(), "html.parser")
            for tag in soup(["script", "style"]):
                tag.decompose()
            return {"text": soup.get_text(separator="\n", strip=True),
                    "source": path, "format": "html"}
        except ImportError:
            with open(path, encoding="utf-8", errors="replace") as f:
                import re
                text = re.sub(r'<[^>]+>', ' ', f.read())
            return {"text": text, "source": path, "format": "html"}

    def batch_parse(self, directory: str, extensions: List[str] = None) -> List[dict]:
        """Parse all documents in a directory."""
        if extensions is None:
            extensions = [".pdf", ".docx", ".txt", ".md", ".rtf", ".pptx", ".epub"]
        results = []
        for path in Path(directory).rglob("*"):
            if path.suffix.lower() in extensions and path.is_file():
                results.append(self.parse(str(path)))
        return results

class DocumentIntelligence:
    """Unified document intelligence interface."""

    def __init__(self):
        self.ocr = ImageOCR()
        self.pdf = PDFExtractor()
        self.parser = DocumentParser()

    def read(self, filepath: str) -> dict:
        """Intelligently read any document or image."""
        path = Path(filepath)
        ext = path.suffix.lower()

        if ext in [".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".webp"]:
            return self.ocr.read_image(filepath)
        else:
            return self.parser.parse(filepath)

    def read_all(self, directory: str) -> List[dict]:
        """Read all documents in a directory."""
        return self.parser.batch_parse(directory)

if __name__ == "__main__":
    di = DocumentIntelligence()
    print("Document Intelligence Module")
    print(f"  Output directory: {OCR_DIR}")
    for name, pkg in [("EasyOCR", "easyocr"), ("pdfplumber", "pdfplumber"),
                       ("docx2txt", "docx2txt"), ("striprtf", "striprtf")]:
        try:
            __import__(pkg)
            print(f"  {name:15s} [OK]")
        except ImportError:
            print(f"  {name:15s} [NOT INSTALLED]")
